from flask import Flask, request, jsonify
from flask_cors import CORS
import os
import time
import PyPDF2
from docx import Document
from pptx import Presentation
import google.generativeai as genai
import json
import re

# --- Flask setup ---
app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# --- Gemini setup ---
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))

# Track last request to avoid hitting free tier quota (2 requests/minute)
last_request_time = 0
REQUEST_DELAY = 35  # seconds


def safe_generate(model, prompt):
    """
    Wrapper to respect free-tier rate limits.
    """
    global last_request_time
    now = time.time()

    if now - last_request_time < REQUEST_DELAY:
        wait_time = REQUEST_DELAY - (now - last_request_time)
        print(f"⏳ Waiting {wait_time:.1f}s before next request...")
        time.sleep(wait_time)

    response = model.generate_content(prompt)
    last_request_time = time.time()
    return response


def extract_text(file_path: str) -> str:
    """
    Extract text from supported file formats: txt, pdf, docx, pptx.
    """
    text = ""
    ext = file_path.lower().split(".")[-1]

    try:
        if ext == "pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"

        elif ext == "docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"

        elif ext == "txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        elif ext in ["ppt", "pptx"]:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        else:
            return None

    except Exception as e:
        print(f"[ERROR] Failed to extract text: {e}")
        return None

    return text.strip()


def clean_json_output(output: str) -> str:
    """
    Remove markdown code fences and extra formatting.
    """
    return re.sub(r"^```[a-zA-Z]*\n|\n```$", "", output.strip())


def generate_quiz(text: str):
    """
    Generate a quiz in JSON format from extracted text.
    """
    if not os.getenv("GEMINI_API_KEY"):
        return {"error": "Missing Gemini API key. Set GEMINI_API_KEY environment variable."}

    prompt = (
        "From the following content, create a multiple-choice quiz in JSON format. "
        "Structure: "
        "[{\"question\": \"...\", \"choices\": [\"...\",\"...\",\"...\"], \"answer\": \"...\"}]. "
        "Generate 20 questions. Return ONLY valid JSON without markdown formatting.\n\n"
        f"Content:\n{text}"
    )

    try:
        model = genai.GenerativeModel("gemini-2.5-flash")
        response = safe_generate(model, prompt)

        quiz_text = ""
        if hasattr(response, "text") and response.text:
            quiz_text = response.text.strip()
        elif hasattr(response, "candidates") and response.candidates:
            quiz_text = response.candidates[0].content.parts[0].text.strip()

        quiz_text = clean_json_output(quiz_text)
        quiz_json = json.loads(quiz_text)

        return quiz_json

    except json.JSONDecodeError:
        return {"error": "Invalid JSON returned by model", "raw": quiz_text}
    except Exception as e:
        print(f"[ERROR] Quiz generation failed: {e}")
        return {"error": str(e)}


def generate_quiz_title(text: str):
    """
    Generate a quiz title and instructions.
    """
    if not os.getenv("GEMINI_API_KEY"):
        return {"error": "Missing Gemini API key. Set GEMINI_API_KEY environment variable."}

    prompt = (
        f"Generate a quiz title and instruction based on this content:\n{text}\n"
        "Structure: "
        "[{\"title\": \"...\", \"instruction\": \"...\"}]. "
        "Return ONLY valid JSON."
    )

    try:
        model = genai.GenerativeModel("gemini-2.5-flash")
        response = safe_generate(model, prompt)

        quiz_text = ""
        if hasattr(response, "text") and response.text:
            quiz_text = response.text.strip()
        elif hasattr(response, "candidates") and response.candidates:
            quiz_text = response.candidates[0].content.parts[0].text.strip()

        quiz_text = clean_json_output(quiz_text)
        quiz_json = json.loads(quiz_text)
        return quiz_json

    except json.JSONDecodeError:
        return {"error": "Invalid JSON returned by model", "raw": quiz_text}
    except Exception as e:
        print(f"[ERROR] Quiz title generation failed: {e}")
        return {"error": str(e)}


@app.route("/upload", methods=["POST"])
def upload_file():
    """
    Upload a file, extract text, and generate quiz + title.
    """
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400

    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(file_path)

    text = extract_text(file_path)
    if not text:
        return jsonify({"error": "Unsupported or unreadable file type"}), 400

    quiz = generate_quiz(text)
    title = generate_quiz_title(text)

    return jsonify({"quiz": quiz, "quiz-title": title})


@app.route("/models", methods=["GET"])
def list_models():
    """
    List available Gemini models.
    """
    try:
        models = genai.list_models()
        available = []
        for m in models:
            available.append({
                "name": m.name,
                "description": getattr(m, "description", ""),
                "supported_methods": getattr(m, "supported_generation_methods", [])
            })
        return jsonify({"models": available})
    except Exception as e:
        return jsonify({"error": str(e)})


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=9090)
